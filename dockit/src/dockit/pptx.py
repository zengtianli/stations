"""PowerPoint processing — font unification, text formatting, table styling.

Pure logic module: bytes in, bytes out. No file paths, no console output.

Usage:
    from dockit.pptx import standardize

    with open("input.pptx", "rb") as f:
        result = standardize(f.read(), font_name="Microsoft YaHei")

    with open("output.pptx", "wb") as f:
        f.write(result.data)
"""

from dataclasses import dataclass, field
from io import BytesIO

from lxml import etree
from pptx import Presentation
from pptx.oxml.ns import qn

from dockit.text import fix_all


@dataclass
class PptxFormatResult:
    """Result of a PowerPoint formatting operation."""

    data: bytes
    stats: dict[str, int] = field(default_factory=dict)


# -- Font helpers (XML-level manipulation) ------------------------------------


def _set_font_for_run(run, font_name: str):
    """Set font for a single run at XML level (latin + ea + cs)."""
    try:
        run.font.name = font_name
        rPr = run._r.get_or_add_rPr()
        for tag in ("a:latin", "a:ea", "a:cs"):
            elem = rPr.find(qn(tag))
            if elem is None:
                elem = etree.SubElement(rPr, qn(tag))
            elem.set("typeface", font_name)
    except Exception:
        pass


def _set_paragraph_default_font(paragraph, font_name: str):
    """Set paragraph default font properties (defRPr)."""
    try:
        pPr = paragraph._p.get_or_add_pPr()
        defRPr = pPr.find(qn("a:defRPr"))
        if defRPr is None:
            defRPr = etree.SubElement(pPr, qn("a:defRPr"))
        for tag in ("a:latin", "a:ea", "a:cs"):
            elem = defRPr.find(qn(tag))
            if elem is None:
                elem = etree.SubElement(defRPr, qn(tag))
            elem.set("typeface", font_name)
    except Exception:
        pass


def _set_endpara_font(paragraph, font_name: str):
    """Set end-of-paragraph run properties font."""
    try:
        endParaRPr = paragraph._p.find(qn("a:endParaRPr"))
        if endParaRPr is not None:
            for tag in ("a:latin", "a:ea", "a:cs"):
                elem = endParaRPr.find(qn(tag))
                if elem is None:
                    elem = etree.SubElement(endParaRPr, qn(tag))
                elem.set("typeface", font_name)
    except Exception:
        pass


# -- Shape processing (recursive for grouped shapes) --------------------------


def _font_process_shape(shape, font_name: str, stats: dict):
    """Process a shape: set font for all text frames and tables."""
    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            _set_paragraph_default_font(paragraph, font_name)
            _set_endpara_font(paragraph, font_name)
            for run in paragraph.runs:
                _set_font_for_run(run, font_name)
                stats["runs"] += 1
        stats["shapes"] += 1

    if shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                if cell.text_frame:
                    for paragraph in cell.text_frame.paragraphs:
                        _set_paragraph_default_font(paragraph, font_name)
                        _set_endpara_font(paragraph, font_name)
                        for run in paragraph.runs:
                            _set_font_for_run(run, font_name)
                            stats["runs"] += 1
                    stats["tables"] += 1

    if hasattr(shape, "shapes"):
        for sub_shape in shape.shapes:
            _font_process_shape(sub_shape, font_name, stats)


def _format_process_shape(shape, stats: dict):
    """Process a shape: fix text formatting (quotes, punctuation, units)."""
    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if run.text:
                    original = run.text
                    fixed, fix_stats, _ = fix_all(original)
                    stats["quotes"] += fix_stats["quotes"]
                    stats["punctuation"] += fix_stats["punctuation"]
                    stats["units"] += fix_stats["units"]
                    if fixed != original:
                        run.text = fixed

    if shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                if cell.text_frame:
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.text:
                                original = run.text
                                fixed, fix_stats, _ = fix_all(original)
                                stats["quotes"] += fix_stats["quotes"]
                                stats["punctuation"] += fix_stats["punctuation"]
                                stats["units"] += fix_stats["units"]
                                if fixed != original:
                                    run.text = fixed

    if hasattr(shape, "shapes"):
        for sub_shape in shape.shapes:
            _format_process_shape(sub_shape, stats)


def _table_style_shape(shape, stats: dict, *, first_row: bool, horz_banding: bool, first_col: bool):
    """Process a shape: set table style options."""
    if shape.has_table:
        try:
            shape.table.first_row = first_row
            shape.table.horz_banding = horz_banding
            shape.table.first_col = first_col
            stats["tables_styled"] += 1
        except Exception:
            pass

    if hasattr(shape, "shapes"):
        for sub_shape in shape.shapes:
            _table_style_shape(sub_shape, stats, first_row=first_row, horz_banding=horz_banding, first_col=first_col)


def _process_all_shapes(prs, processor, stats: dict, *, include_masters: bool = True):
    """Apply a processor function to all shapes in the presentation."""
    if include_masters:
        for slide_master in prs.slide_masters:
            for shape in slide_master.shapes:
                processor(shape, stats)
            for layout in slide_master.slide_layouts:
                for shape in layout.shapes:
                    processor(shape, stats)

    for slide in prs.slides:
        for shape in slide.shapes:
            processor(shape, stats)


# -- Public API ----------------------------------------------------------------


def format_text(
    pptx_bytes: bytes,
    *,
    include_masters: bool = True,
) -> PptxFormatResult:
    """Fix text formatting in a PowerPoint file (quotes, punctuation, units).

    Args:
        pptx_bytes: Raw bytes of a .pptx file.
        include_masters: Whether to process slide masters and layouts.

    Returns:
        PptxFormatResult with processed bytes and stats.
    """
    prs = Presentation(BytesIO(pptx_bytes))
    stats = {"quotes": 0, "punctuation": 0, "units": 0}

    _process_all_shapes(prs, _format_process_shape, stats, include_masters=include_masters)

    buf = BytesIO()
    prs.save(buf)
    return PptxFormatResult(data=buf.getvalue(), stats=stats)


def set_font(
    pptx_bytes: bytes,
    font_name: str = "Microsoft YaHei",
    *,
    include_masters: bool = True,
) -> PptxFormatResult:
    """Unify all fonts in a PowerPoint file.

    Args:
        pptx_bytes: Raw bytes of a .pptx file.
        font_name: Target font name.
        include_masters: Whether to process slide masters and layouts.

    Returns:
        PptxFormatResult with processed bytes and stats.
    """
    prs = Presentation(BytesIO(pptx_bytes))
    stats = {"runs": 0, "shapes": 0, "tables": 0}

    def processor(shape, s):
        _font_process_shape(shape, font_name, s)

    _process_all_shapes(prs, processor, stats, include_masters=include_masters)

    buf = BytesIO()
    prs.save(buf)
    return PptxFormatResult(data=buf.getvalue(), stats=stats)


def set_table_style(
    pptx_bytes: bytes,
    *,
    first_row: bool = True,
    horz_banding: bool = True,
    first_col: bool = True,
) -> PptxFormatResult:
    """Set table style options across all slides.

    Args:
        pptx_bytes: Raw bytes of a .pptx file.
        first_row: Enable header row styling.
        horz_banding: Enable banded rows.
        first_col: Enable first column styling.

    Returns:
        PptxFormatResult with processed bytes and stats.
    """
    prs = Presentation(BytesIO(pptx_bytes))
    stats = {"tables_styled": 0}

    def processor(shape, s):
        _table_style_shape(shape, s, first_row=first_row, horz_banding=horz_banding, first_col=first_col)

    _process_all_shapes(prs, processor, stats, include_masters=False)

    buf = BytesIO()
    prs.save(buf)
    return PptxFormatResult(data=buf.getvalue(), stats=stats)


def standardize(
    pptx_bytes: bytes,
    font_name: str = "Microsoft YaHei",
) -> PptxFormatResult:
    """Apply all standardization: format_text -> set_font -> set_table_style.

    Args:
        pptx_bytes: Raw bytes of a .pptx file.
        font_name: Target font name for unification.

    Returns:
        PptxFormatResult with processed bytes and combined stats.
    """
    combined_stats: dict[str, int] = {}

    # Step 1: Format text
    result = format_text(pptx_bytes)
    combined_stats.update(result.stats)

    # Step 2: Set font
    result = set_font(result.data, font_name)
    combined_stats.update(result.stats)

    # Step 3: Set table style
    result = set_table_style(result.data)
    combined_stats.update(result.stats)

    return PptxFormatResult(data=result.data, stats=combined_stats)


def to_markdown(pptx_bytes: bytes) -> str:
    """Convert a PowerPoint file to Markdown text.

    Extracts text from all slides including shapes and speaker notes.

    Args:
        pptx_bytes: Raw bytes of a .pptx file.

    Returns:
        Markdown-formatted string with slide content.
    """
    prs = Presentation(BytesIO(pptx_bytes))
    parts = []

    for i, slide in enumerate(prs.slides, 1):
        parts.append(f"## Slide {i}\n")

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(shape.text + "\n")

        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text
            if notes:
                parts.append(f"### Speaker Notes\n\n{notes}\n")

        parts.append("---\n")

    return "\n".join(parts)
